"""
PPTX service for PowerPoint presentation manipulation using python-pptx
"""
import re
from pathlib import Path
from typing import List, Optional, Set
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Pt
from pptx.dml.color import RGBColor

from app.models.schemas import (
    VariableInfo,
    TemplateVariables,
    TextFormatting
)
from app.models.enums import TextAlignment, VerticalAlignment
from app.services.file_service import FileService


class PPTXService:
    """Service for PowerPoint presentation operations using curly brace variables"""
    
    def __init__(self, file_service: FileService):
        """
        Initialize PPTX service
        
        Args:
            file_service: File service instance
        """
        self.file_service = file_service
        self.var_regex = re.compile(r"\{\{(.*?)\}\}")
    
    def _get_alt_text(self, shape) -> Optional[str]:
        """Helper to extract Alt Text from a shape's XML"""
        try:
            nvPr = shape._element.nvXxPr
            cNvPr = nvPr.find('p:cNvPr', namespaces=nvPr.nsmap)
            if cNvPr is None:
                cNvPr = nvPr.find('*/p:cNvPr', namespaces=nvPr.nsmap)
            
            if cNvPr is not None:
                return cNvPr.get("descr")
        except Exception:
            pass
        return None

    def get_template_variables(self, template_id: str) -> TemplateVariables:
        """
        Get all {{variable}} patterns from a template
        
        Args:
            template_id: Template ID
            
        Returns:
            TemplateVariables object
        """
        template_path = self.file_service.get_template_path(template_id)
        
        try:
            prs = Presentation(str(template_path))
        except Exception as e:
            raise Exception(f"Failed to load template: {str(e)}")
        
        variables: List[VariableInfo] = []
        seen_vars: Set[str] = set()

        for slide_idx, slide in enumerate(prs.slides):
            for shape in slide.shapes:
                # 1. Check Text variables
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        matches = self.var_regex.findall(paragraph.text)
                        for match in matches:
                            var_key = f"text:{match}:{slide_idx}"
                            if var_key not in seen_vars:
                                variables.append(VariableInfo(
                                    name=match,
                                    type="text",
                                    slide_index=slide_idx
                                ))
                                seen_vars.add(var_key)

                # 2. Check Image variables in Alt Text
                alt_text = self._get_alt_text(shape)
                if alt_text:
                    matches = self.var_regex.findall(alt_text)
                    for match in matches:
                        # Normalize image variables (remove 'image:' prefix if present)
                        clean_match = match.replace("image:", "")
                        var_key = f"image:{clean_match}:{slide_idx}"
                        if var_key not in seen_vars:
                            variables.append(VariableInfo(
                                name=clean_match,
                                type="image",
                                slide_index=slide_idx
                            ))
                            seen_vars.add(var_key)
        
        return TemplateVariables(
            template_id=template_id,
            variables=variables
        )
    
    def create_presentation(self, template_id: str, presentation_id: str) -> str:
        """Create a new presentation from a template"""
        template_path = self.file_service.get_template_path(template_id)
        
        try:
            prs = Presentation(str(template_path))
        except Exception as e:
            raise Exception(f"Failed to load template: {str(e)}")
        
        output_path = self.file_service.create_presentation_path(presentation_id)
        
        try:
            prs.save(str(output_path))
        except Exception as e:
            raise Exception(f"Failed to save presentation: {str(e)}")
        
        return str(output_path)
    
    def insert_text(
        self,
        presentation_id: str,
        variable_name: str,
        text: str,
        formatting: Optional[TextFormatting] = None
    ) -> bool:
        """
        Global search and replace for {{variable_name}}
        """
        presentation_path = self.file_service.get_presentation_path(presentation_id)
        
        try:
            prs = Presentation(str(presentation_path))
        except Exception as e:
            raise Exception(f"Failed to load presentation: {str(e)}")
        
        target = "{{" + variable_name + "}}"
        replaced = False
        
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    # Check if the target is in the text frame at all
                    if target in shape.text_frame.text:
                        replaced = True
                        # Replace in each paragraph to preserve formatting as much as possible
                        for paragraph in shape.text_frame.paragraphs:
                            if target in paragraph.text:
                                # If the paragraph ONLY contains the variable, we can apply global formatting
                                if paragraph.text.strip() == target:
                                    paragraph.text = text
                                    if formatting:
                                        self._apply_paragraph_formatting(paragraph, formatting)
                                else:
                                    # Mixed text: simple replacement
                                    # Note: this might lose some run-level formatting if the variable spans runs
                                    paragraph.text = paragraph.text.replace(target, text)
        
        if not replaced:
            # We don't raise exception for text as there might be many variables in a slide
            pass
        
        try:
            prs.save(str(presentation_path))
        except Exception as e:
            raise Exception(f"Failed to save presentation: {str(e)}")
        
        return True

    def insert_image(
        self,
        presentation_id: str,
        variable_name: str,
        image_path: str
    ) -> bool:
        """
        Replace image by finding {{variable_name}} or {{image:variable_name}} in Alt Text
        """
        presentation_path = self.file_service.get_presentation_path(presentation_id)
        
        try:
            prs = Presentation(str(presentation_path))
        except Exception as e:
            raise Exception(f"Failed to load presentation: {str(e)}")
        
        image_replaced = False
        
        # We look for {{var}} or {{image:var}}
        patterns = ["{{" + variable_name + "}}", "{{image:" + variable_name + "}}"]
        
        for slide in prs.slides:
            for shape in list(slide.shapes):
                alt_text = self._get_alt_text(shape)
                if alt_text and any(p == alt_text.strip() for p in patterns):
                    image_replaced = True
                    
                    # Geometry
                    left, top = shape.left, shape.top
                    width, height = shape.width, shape.height
                    
                    try:
                        slide.shapes.add_picture(image_path, left, top, width, height)
                        # Remove original
                        sp = shape._element
                        sp.getparent().remove(sp)
                    except Exception as e:
                        raise Exception(f"Failed to replace image shape: {str(e)}")
                    
                    # Continue searching (could be multiple instances)
            
        if not image_replaced:
            raise Exception(f"No image variable found with Alt Text '{{{{{variable_name}}}}}'")
        
        try:
            prs.save(str(presentation_path))
        except Exception as e:
            raise Exception(f"Failed to save presentation: {str(e)}")
        
        return True

    def _apply_paragraph_formatting(self, paragraph, formatting: TextFormatting):
        """Apply formatting to a paragraph and its runs"""
        # Alignment
        if formatting.alignment:
            alignment_mapping = {
                TextAlignment.LEFT: PP_ALIGN.LEFT,
                TextAlignment.CENTER: PP_ALIGN.CENTER,
                TextAlignment.RIGHT: PP_ALIGN.RIGHT,
                TextAlignment.JUSTIFY: PP_ALIGN.JUSTIFY,
                TextAlignment.DISTRIBUTE: PP_ALIGN.DISTRIBUTE
            }
            paragraph.alignment = alignment_mapping.get(formatting.alignment, PP_ALIGN.LEFT)
        
        # Font formatting (applied to all runs)
        for run in paragraph.runs:
            font = run.font
            if formatting.font_name:
                font.name = formatting.font_name
            if formatting.font_size:
                font.size = Pt(formatting.font_size)
            if formatting.bold is not None:
                font.bold = formatting.bold
            if formatting.italic is not None:
                font.italic = formatting.italic
            if formatting.underline is not None:
                font.underline = formatting.underline
            if formatting.color:
                hex_color = formatting.color.lstrip('#')
                r, g, b = int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16)
                font.color.rgb = RGBColor(r, g, b)
